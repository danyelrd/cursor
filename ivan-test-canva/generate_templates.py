#!/usr/bin/env python3
"""Generate Ivan Test Canva-editable SVG + PowerPoint templates."""

from pathlib import Path

import svgwrite
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

OUT = Path(__file__).resolve().parent
SVG_DIR = OUT / "svg"
PPTX_DIR = OUT / "pptx"

# 1080px square at 96dpi
SLIDE_EMU = Emu(10287000)

GOLD = "#C9A227"
CREAM = "#F5F0E8"
CREAM_DARK = "#E8DFD0"
BROWN = "#3D2914"
GREEN = "#2E7D32"
WHITE = "#FFFFFF"


def rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_rect(slide, left, top, width, height, fill_hex, line_hex=None, line_width=0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    if line_hex:
        shape.line.color.rgb = rgb(line_hex)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    size_pt,
    color_hex,
    bold=False,
    align=PP_ALIGN.LEFT,
    font_name="Arial",
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = rgb(color_hex)
    run.font.name = font_name
    return box


def add_placeholder(slide, left, top, width, height, label):
    shape = add_rect(slide, left, top, width, height, CREAM_DARK, GOLD, 2)
    shape.fill.transparency = 0.15
    add_textbox(
        slide,
        left,
        top,
        width,
        height,
        label,
        14,
        BROWN,
        align=PP_ALIGN.CENTER,
    )
    return shape


def new_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_EMU
    prs.slide_height = SLIDE_EMU
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def save_pptx(prs, name: str):
    path = PPTX_DIR / f"{name}.pptx"
    prs.save(path)
    return path


def svg_text(
    dwg,
    x,
    y,
    text,
    size,
    fill,
    weight="normal",
    anchor="start",
    font_family="Arial, Helvetica, sans-serif",
):
    dwg.add(
        dwg.text(
            text,
            insert=(x, y),
            fill=fill,
            font_size=size,
            font_family=font_family,
            font_weight=weight,
            text_anchor=anchor,
        )
    )


def svg_placeholder(dwg, x, y, w, h, label):
    dwg.add(
        dwg.rect(
            insert=(x, y),
            size=(w, h),
            fill=CREAM_DARK,
            stroke=GOLD,
            stroke_width=2,
            rx=8,
            ry=8,
        )
    )
    svg_text(dwg, x + w / 2, y + h / 2 + 6, label, 18, BROWN, anchor="middle")


def build_slide_1_pptx():
    prs = new_presentation()
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_EMU, SLIDE_EMU, CREAM)
    add_textbox(
        slide,
        Emu(400000),
        Emu(350000),
        Emu(9500000),
        Emu(900000),
        "IVAN TEST",
        22,
        GOLD,
        bold=True,
        font_name="Georgia",
    )
    add_textbox(
        slide,
        Emu(400000),
        Emu(1100000),
        Emu(9500000),
        Emu(1200000),
        "CYSTEINE SYSTEM",
        44,
        GOLD,
        bold=True,
        align=PP_ALIGN.CENTER,
        font_name="Georgia",
    )
    add_placeholder(
        slide,
        Emu(1200000),
        Emu(2600000),
        Emu(4200000),
        Emu(4200000),
        "Replace:\nHair strand macro\n+ gold droplets",
    )
    add_placeholder(
        slide,
        Emu(5600000),
        Emu(3000000),
        Emu(3200000),
        Emu(5200000),
        "Replace:\nIvan Test\nproduct bottle",
    )
    add_textbox(
        slide,
        Emu(700000),
        Emu(7200000),
        Emu(8800000),
        Emu(1100000),
        "Penetrates deep into hair fibers for better results.",
        24,
        BROWN,
        align=PP_ALIGN.CENTER,
    )
    return save_pptx(prs, "01-ivan-test-cysteine")


def build_slide_2_pptx():
    prs = new_presentation()
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_EMU, SLIDE_EMU, CREAM)
    add_textbox(
        slide,
        Emu(500000),
        Emu(500000),
        Emu(2000000),
        Emu(600000),
        "Ivan Test",
        20,
        BROWN,
        bold=True,
        font_name="Georgia",
    )
    add_textbox(
        slide,
        Emu(500000),
        Emu(1400000),
        Emu(9300000),
        Emu(2000000),
        "Bakit importante ang FDA notified products?",
        36,
        BROWN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Emu(5500000),
        Emu(2500000),
        Emu(4200000),
        Emu(500000),
        "FDA notified",
        28,
        GREEN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_placeholder(
        slide,
        Emu(600000),
        Emu(3600000),
        Emu(9000000),
        Emu(3200000),
        "Replace: 4–5 Ivan Test bottles\n(same backdrop + lighting)",
    )
    add_textbox(
        slide,
        Emu(700000),
        Emu(7200000),
        Emu(8800000),
        Emu(1400000),
        "Maingat na formula • Salon-trusted • Ivan Test quality",
        20,
        BROWN,
        align=PP_ALIGN.CENTER,
    )
    return save_pptx(prs, "02-ivan-test-fda-tagalog")


def build_slide_3_pptx():
    prs = new_presentation()
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_EMU, SLIDE_EMU, CREAM_DARK)
    add_textbox(
        slide,
        Emu(400000),
        Emu(400000),
        Emu(9400000),
        Emu(1100000),
        "Ivan Test Argan Oil",
        40,
        GOLD,
        bold=True,
        align=PP_ALIGN.CENTER,
        font_name="Georgia",
    )
    add_placeholder(
        slide,
        Emu(200000),
        Emu(1800000),
        Emu(5000000),
        Emu(6200000),
        "Replace:\nModel photo\n(glossy hair)",
    )
    add_placeholder(
        slide,
        Emu(5400000),
        Emu(2800000),
        Emu(4600000),
        Emu(5200000),
        "Replace:\nIvan Test bottle\n+ soft shadow",
    )
    add_rect(slide, Emu(5200000), Emu(1900000), Emu(4200000), Emu(700000), CREAM, GOLD, 1)
    add_textbox(
        slide,
        Emu(5300000),
        Emu(1950000),
        Emu(4000000),
        Emu(600000),
        "shine and smoothness",
        22,
        BROWN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Emu(500000),
        Emu(8200000),
        Emu(9200000),
        Emu(800000),
        "Keratin + Argan blend for salon-smooth finish",
        18,
        BROWN,
        align=PP_ALIGN.CENTER,
    )
    return save_pptx(prs, "03-ivan-test-argan")


def build_slide_4_pptx():
    prs = new_presentation()
    slide = blank_slide(prs)
    add_placeholder(
        slide,
        Emu(0),
        Emu(0),
        SLIDE_EMU,
        SLIDE_EMU,
        "Replace:\nSalon photo\n(stylist + client)",
    )
    add_rect(slide, Emu(0), Emu(6200000), SLIDE_EMU, Emu(4067000), BROWN)
    shape = slide.shapes[-1]
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(BROWN)
    shape.fill.transparency = 0.35
    add_textbox(
        slide,
        Emu(600000),
        Emu(6600000),
        Emu(9000000),
        Emu(2000000),
        "Many salon owners recommend Ivan Test for their clients.",
        34,
        WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        font_name="Georgia",
    )
    add_placeholder(
        slide,
        Emu(1200000),
        Emu(8800000),
        Emu(7600000),
        Emu(1100000),
        "Replace: small Ivan Test product strip",
    )
    return save_pptx(prs, "04-ivan-test-salon-proof")


def build_svgs():
    specs = [
        (
            "01-ivan-test-cysteine.svg",
            lambda dwg: (
                dwg.add(dwg.rect(insert=(0, 0), size=("1080", "1080"), fill=CREAM)),
                svg_text(dwg, 540, 120, "IVAN TEST", 28, GOLD, weight="bold", anchor="middle", font_family="Georgia, serif"),
                svg_text(dwg, 540, 175, "CYSTEINE SYSTEM", 48, GOLD, weight="bold", anchor="middle", font_family="Georgia, serif"),
                svg_placeholder(dwg, 100, 280, 480, 480, "Hair strand + droplets"),
                svg_placeholder(dwg, 620, 320, 360, 520, "Product bottle"),
                svg_text(
                    dwg,
                    540,
                    920,
                    "Penetrates deep into hair fibers for better results.",
                    26,
                    BROWN,
                    anchor="middle",
                ),
            ),
        ),
        (
            "02-ivan-test-fda-tagalog.svg",
            lambda dwg: (
                dwg.add(dwg.rect(insert=(0, 0), size=("1080", "1080"), fill=CREAM)),
                svg_text(dwg, 48, 56, "Ivan Test", 24, BROWN, weight="bold", font_family="Georgia, serif"),
                svg_text(
                    dwg,
                    540,
                    150,
                    "Bakit importante ang FDA notified products?",
                    34,
                    BROWN,
                    weight="bold",
                    anchor="middle",
                ),
                svg_text(dwg, 700, 210, "FDA notified", 30, GREEN, weight="bold"),
                svg_placeholder(dwg, 60, 300, 960, 360, "Product bottle lineup"),
                svg_text(
                    dwg,
                    540,
                    820,
                    "Maingat na formula • Salon-trusted • Ivan Test quality",
                    22,
                    BROWN,
                    anchor="middle",
                ),
            ),
        ),
        (
            "03-ivan-test-argan.svg",
            lambda dwg: (
                dwg.add(dwg.rect(insert=(0, 0), size=("1080", "1080"), fill=CREAM_DARK)),
                svg_text(dwg, 540, 90, "Ivan Test Argan Oil", 44, GOLD, weight="bold", anchor="middle", font_family="Georgia, serif"),
                svg_placeholder(dwg, 40, 170, 500, 620, "Model photo"),
                svg_placeholder(dwg, 560, 280, 460, 500, "Product bottle"),
                dwg.add(
                    dwg.rect(insert=(560, 180), size=(420, 70), fill=CREAM, stroke=GOLD, rx=12, ry=12)
                ),
                svg_text(dwg, 770, 225, "shine and smoothness", 22, BROWN, weight="bold", anchor="middle"),
                svg_text(
                    dwg,
                    540,
                    1000,
                    "Keratin + Argan blend for salon-smooth finish",
                    20,
                    BROWN,
                    anchor="middle",
                ),
            ),
        ),
        (
            "04-ivan-test-salon-proof.svg",
            lambda dwg: (
                svg_placeholder(dwg, 0, 0, 1080, 1080, "Salon background photo"),
                dwg.add(
                    dwg.rect(insert=(0, 620), size=(1080, 460), fill=BROWN, opacity=0.55)
                ),
                svg_text(
                    dwg,
                    540,
                    760,
                    "Many salon owners recommend",
                    32,
                    WHITE,
                    weight="bold",
                    anchor="middle",
                    font_family="Georgia, serif",
                ),
                svg_text(
                    dwg,
                    540,
                    810,
                    "Ivan Test for their clients.",
                    36,
                    GOLD,
                    weight="bold",
                    anchor="middle",
                    font_family="Georgia, serif",
                ),
                svg_placeholder(dwg, 140, 900, 800, 120, "Product strip"),
            ),
        ),
    ]
    paths = []
    for filename, draw_fn in specs:
        dwg = svgwrite.Drawing(
            str(SVG_DIR / filename),
            size=("1080px", "1080px"),
            viewBox="0 0 1080 1080",
        )
        draw_fn(dwg)
        dwg.save()
        paths.append(SVG_DIR / filename)
    return paths


def main():
    SVG_DIR.mkdir(parents=True, exist_ok=True)
    PPTX_DIR.mkdir(parents=True, exist_ok=True)
    pptx_paths = [
        build_slide_1_pptx(),
        build_slide_2_pptx(),
        build_slide_3_pptx(),
        build_slide_4_pptx(),
    ]
    svg_paths = build_svgs()
    print("Created:")
    for p in pptx_paths + svg_paths:
        print(" ", p)


if __name__ == "__main__":
    main()
